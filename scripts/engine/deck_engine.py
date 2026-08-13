#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
deck_engine.py —— 20 页 PPT 共享渲染引擎
==========================================
从 build_geo_deck.py 抽取的渲染逻辑，参数化为：
    build_deck(config: dict, palette: dict, output_path: str)

- config：case 的 CONFIG dict（品牌 / AIVO / 曝光 / 频次 / 竞品 / 话题词 / KPI / 增强变量）
- palette：来自 engine/palette.py 的预设或自定义 palette（配色/字体作为数据传入，不硬编码）
- 20 页结构、坐标、尺寸 1:1 保留自原 build_geo_deck.py，不破坏 OOB 检查。

color_map 由 palette 动态构建：slide 绘制代码用 BLUE/NAVY/CYAN/... 逻辑色名，
其实际 RGB 由当前 case 的 palette 决定。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

from .palette import build_rgb, get_palette


# ------------------------------------------------------------------
# CONFIG 容量校验（版面硬约束；超出必然 OOB，宁可提前报错）
# 由 build_deck 入口调用，扫描 config 内的列表型字段
# ------------------------------------------------------------------
CONFIG_LIMITS = [
    ("PLATFORM_INFLUENCE", 7),   # 含表头，第 5 页
    ("FREQ_ROWS", 6),            # 第 7 页
    ("SOURCE_TABLE", 4),         # 第 8 页
    ("INFRA_CARDS", 3),          # 第 11 页
    ("CURRENT_FACTS", 5),        # 第 13 页
    ("ECO_MAP", 5),              # 第 14 页上表
    ("COMPETITORS", 4),          # 第 14 页下表
    ("GAP_ROWS", 6),             # 第 16 页
    ("CAPABILITY_MAP", 6),       # 第 17 页
    ("TOPIC_WORDS", 5),          # 第 18 页话题卡
    ("KPI_ROWS", 5),             # 含表头，第 18 页
    ("KEY_FINDINGS", 5),         # 第 2 页
    ("STAT_CARDS", 3),           # 第 4 页
]


def _validate_config(config):
    for _name, _max in CONFIG_LIMITS:
        _seq = config.get(_name)
        if _seq is not None and len(_seq) > _max:
            raise SystemExit(
                f"[CONFIG 容量超限] {_name} 有 {len(_seq)} 项，上限 {_max} 项。"
                f"超出会导致该页元素越界（OOB），请合并或精简后重跑。")


def build_deck(config, palette, output_path):
    """根据 case config + palette 生成 20 页 PPT，返回输出路径。"""
    _validate_config(config)

    pal = get_palette(palette)
    C = build_rgb(pal)                 # {逻辑色名: RGBColor}
    CM = dict(C)                       # COLOR_MAP：逻辑色名 → RGBColor
    FONT = pal["font"]
    FNUM = pal["font_numeral"]
    chart_series = pal.get("chart_series", ["#2A5BEA", "#9DB6F5"])
    _cs0 = RGBColor(int(chart_series[0].lstrip("#")[0:2], 16),
                    int(chart_series[0].lstrip("#")[2:4], 16),
                    int(chart_series[0].lstrip("#")[4:6], 16))
    _cs1 = RGBColor(int(chart_series[1].lstrip("#")[0:2], 16),
                    int(chart_series[1].lstrip("#")[2:4], 16),
                    int(chart_series[1].lstrip("#")[4:6], 16))

    EMU = 914400
    SW, SH = Inches(13.333), Inches(7.5)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # ---- 从 config 取变量（含默认值，增强变量未设置时原样输出）----
    G = dict(config)
    def V(k, default=None): return G.get(k, default)

    BRAND_CN = V("BRAND_CN"); BRAND_EN = V("BRAND_EN")
    CLIENT_DESC = V("CLIENT_DESC"); PRODUCT_TYPE = V("PRODUCT_TYPE")
    CONTACT = V("CONTACT", [])
    DIAG = V("DIAG", "微盟星启 GEO 診斷模型估算")
    AIVO_TOTAL = V("AIVO_TOTAL"); AIVO_RATE = V("AIVO_RATE")
    AIVO_VIS, AIVO_INFRA, AIVO_COMP, AIVO_SENT = V("AIVO_VIS", 0), V("AIVO_INFRA", 0), V("AIVO_COMP", 0), V("AIVO_SENT", 0)
    AIVO_BENCH = V("AIVO_BENCH", [55, 68, 60, 75])
    CITE_RATE = V("CITE_RATE"); CITED = V("CITED"); TOTAL_SCEN = V("TOTAL_SCEN"); GAP_PCT = V("GAP_PCT")
    PLATFORM_INFLUENCE = V("PLATFORM_INFLUENCE"); INFLUENCE_NOTE = V("INFLUENCE_NOTE")
    FREQ_ROWS = V("FREQ_ROWS"); FREQ_HL = V("FREQ_HL")
    COMPETITORS = V("COMPETITORS"); ECO_MAP = V("ECO_MAP")
    COMPETITOR_TACTICS = V("COMPETITOR_TACTICS"); COMPETITOR_HL = V("COMPETITOR_HL")
    STAT_CARDS = V("STAT_CARDS"); INDUSTRY_NOTE = V("INDUSTRY_NOTE")
    INFRA_CARDS = V("INFRA_CARDS")
    KPI_ROWS = V("KPI_ROWS"); QA_GRAPH = V("QA_GRAPH")
    SOURCE_OFFICIAL = V("SOURCE_OFFICIAL"); SOURCE_TABLE = V("SOURCE_TABLE"); SOURCE_FOOTNOTE = V("SOURCE_FOOTNOTE")
    TOPIC_WORDS = V("TOPIC_WORDS"); KEY_FINDINGS = V("KEY_FINDINGS")
    SENTIMENT_BULLETS = V("SENTIMENT_BULLETS"); SENTIMENT_ACTIONS = V("SENTIMENT_ACTIONS")
    CURRENT_FACTS = V("CURRENT_FACTS"); GAP_ROWS = V("GAP_ROWS")
    CAPABILITY_MAP = V("CAPABILITY_MAP"); ROADMAP_TABLE = V("ROADMAP_TABLE")
    # 增强变量（可选）
    VIS_DUAL = V("VIS_DUAL"); RANK_POOL = V("RANK_POOL")
    SOURCE_CITE = V("SOURCE_CITE"); COMPLIANCE = V("COMPLIANCE", {"applicable": False})

    # ---- 工具函数 ----
    def slide(): return prs.slides.add_slide(blank)

    def bg(s, color): 
        s.background.fill.solid(); s.background.fill.fore_color.rgb = color

    def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
        sp = s.shapes.add_shape(shape, l, t, w, h)
        if fill is None: sp.fill.background()
        else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
        if line is None: sp.line.fill.background()
        else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
        sp.shadow.inherit = False
        return sp

    def txt(s, l, t, w, h, text, size=14, color=None, bold=False, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=None, italic=False):
        color = C["INK"] if color is None else color
        font = FONT if font is None else font
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = Pt(2); tf.margin_right = Pt(2)
        tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.name = font; run.font.color.rgb = color
        return tb

    def bullets(s, l, t, w, h, items, size=14, color=None, gap=6, bullet="•  "):
        color = C["INK"] if color is None else color
        tb = s.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap); p.alignment = PP_ALIGN.LEFT
            run = p.add_run(); run.text = bullet + it
            run.font.size = Pt(size); run.font.name = FONT; run.font.color.rgb = color
        return tb

    def header(s, title, sub=None):
        rect(s, Inches(0.5), Inches(0.42), Inches(0.16), Inches(0.62), fill=C["BLUE"])
        txt(s, Inches(0.8), Inches(0.38), Inches(11.5), Inches(0.7), title,
            size=30, color=C["INK"], bold=True)
        if sub:
            txt(s, Inches(0.82), Inches(1.05), Inches(11.5), Inches(0.45), sub, size=14, color=C["GRAY"])
        txt(s, Inches(0.8), Inches(7.02), Inches(11.5), Inches(0.35),
            "微盟星启GEO运营团队", size=9, color=C["GRAY"])

    def table(s, x, y, w, h, data, col_widths, fs=11, rh=0.34, header_fill=None):
        header_fill = C["BLUE"] if header_fill is None else header_fill
        nrows = len(data); ncols = len(data[0])
        gtbl = s.shapes.add_table(nrows, ncols, x, y, w, Inches(rh * nrows + 0.1))
        tbl = gtbl.table; tbl.first_row = True; tbl.horz_banding = False
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            tbl.columns[j].width = Emu(int(w * cw / total))
        for i, row in enumerate(data):
            tbl.rows[i].height = Inches(rh + (0.06 if i == 0 else 0))
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.margin_left = Inches(0.07); cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill if i == 0 else (C["WHITE"] if i % 2 == 1 else C["CLOUD"])
                tf = cell.text_frame; tf.word_wrap = True
                p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
                run = p.add_run(); run.text = str(val)
                run.font.size = Pt(fs); run.font.name = FONT
                run.font.bold = (i == 0)
                run.font.color.rgb = C["WHITE"] if i == 0 else C["INK"]
        return gtbl

    # ================================================================
    # Slide 1 — 封面
    # ================================================================
    s = slide(); bg(s, C["NAVY"])
    rect(s, 0, SH - Inches(2.2), SW, Inches(2.2), fill=C["BLUE"])
    rect(s, 0, 0, Inches(0.25), SH, fill=C["CYAN"])
    rect(s, Inches(0.9), Inches(0.8), Inches(0.5), Inches(0.12), fill=C["CYAN"])
    txt(s, Inches(0.9), Inches(1.0), Inches(6), Inches(0.5), "微盟星启 GEO", size=16, color=C["WHITE"], bold=True)
    txt(s, Inches(0.85), Inches(2.5), Inches(11.6), Inches(1.3), BRAND_CN, size=42, color=C["WHITE"], bold=True)
    txt(s, Inches(0.9), Inches(3.75), Inches(11.6), Inches(0.8), "微盟星启 GEO 優化方案", size=30, color=C["CYAN"], bold=True)
    txt(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.5), "國內版 6 平臺（主） · 海外版 5 平臺（輔）雙版本佈局", size=15, color=C["WHITE"])
    txt(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.5), "讓你的品牌，在 AI 推薦中被看見", size=18, color=C["WHITE"])
    txt(s, Inches(10.6), Inches(5.55), Inches(2.4), Inches(0.5), "微盟 Weimob", size=14, color=C["WHITE"], align=PP_ALIGN.RIGHT)

    # ================================================================
    # Slide 2 — 執行摘要
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "執行摘要 · 結論先行", "基於" + DIAG)
    rect(s, Inches(0.8), Inches(1.5), Inches(3.4), Inches(2.4), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.8), Inches(1.7), Inches(3.4), Inches(0.5), "AIVO 總評", size=16, color=C["CYAN"], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(2.15), Inches(3.4), Inches(1.1), str(AIVO_TOTAL), size=66, color=C["WHITE"], bold=True, align=PP_ALIGN.CENTER, font=FNUM)
    txt(s, Inches(0.8), Inches(3.25), Inches(3.4), Inches(0.5), AIVO_RATE + "（≥60）", size=16, color=C["WHITE"], align=PP_ALIGN.CENTER, bold=True)
    dims = [("AI 搜索可見性", str(AIVO_VIS), C["AMBER"]), ("基建完善度", str(AIVO_INFRA), C["BLUE"]),
            ("競爭優勢", str(AIVO_COMP), C["BLUE"]), ("輿情健康度", str(AIVO_SENT), C["GREEN"])]
    x = Inches(4.5)
    for name, sc, col in dims:
        rect(s, x, Inches(1.5), Inches(2.0), Inches(1.1), fill=C["WHITE"], line=col, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, Inches(1.62), Inches(2.0), Inches(0.5), sc, size=30, color=col, bold=True, align=PP_ALIGN.CENTER, font=FNUM)
        txt(s, x, Inches(2.18), Inches(2.0), Inches(0.4), name, size=12, color=C["INK"], align=PP_ALIGN.CENTER)
        x += Inches(2.12)
    txt(s, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4), "關鍵發現", size=18, color=C["BLUE"], bold=True)
    bullets(s, Inches(0.8), Inches(4.7), Inches(11.8), Inches(2.2), KEY_FINDINGS, size=13.5, gap=7)

    # ================================================================
    # Slide 3 — 方法論
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "方法論 · GEO 診斷 + 數據追蹤", "AIVO 四維評分 · 4 階段流水線 · 平臺雙版本")
    rect(s, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.0), fill=C["WHITE"], line=C["BLUE"], line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(1.05), Inches(1.7), Inches(5.2), Inches(0.5), "GEO 診斷引擎（AIVO 四維）", size=17, color=C["BLUE"], bold=True)
    bullets(s, Inches(1.05), Inches(2.35), Inches(5.3), Inches(2.0), [
        "AI 搜索可見性：常見問答覆蓋率 / 被引用率 / 推薦位佔比",
        "基建完善度：官網 / 自媒體矩陣 / 權威媒體 / 結構化數據",
        "競爭優勢：相對競品在 AI 生態的位置與差異化",
        "輿情健康度：正面 / 負面情緒與信息來源追蹤",
    ], size=13, gap=7)
    txt(s, Inches(1.05), Inches(4.55), Inches(5.2), Inches(0.5), "4 階段流水線", size=17, color=C["BLUE"], bold=True)
    bullets(s, Inches(1.05), Inches(5.15), Inches(5.3), Inches(1.3), [
        "① 基礎調研（用戶畫像 / 基建 / 競品）",
        "② 收錄 + 可見性（並行）  ③ 輿情分析（並行）",
        "④ 評分 + 建議",
    ], size=13, gap=6)
    rect(s, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.0), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(7.05), Inches(1.7), Inches(5.2), Inches(0.5), "微盟星启 GEO 全鏈方案", size=17, color=C["CYAN"], bold=True)
    bullets(s, Inches(7.05), Inches(2.35), Inches(5.3), Inches(1.7), [
        "智能診斷 → 內容優化 → 全域分發 → 持續監測",
        "四大能力模塊：AI 可見性監測 / 品牌輿情指數監測",
        "內容創作與改造優化 / 智能媒體匹配與發布",
    ], size=13, color=C["WHITE"], gap=7)
    txt(s, Inches(7.05), Inches(4.35), Inches(5.2), Inches(0.5), "平臺雙版本（固定）", size=17, color=C["CYAN"], bold=True)
    txt(s, Inches(7.05), Inches(4.95), Inches(5.3), Inches(0.5), "國內版 6（主）：", size=13, color=C["WHITE"], bold=True)
    txt(s, Inches(7.05), Inches(5.35), Inches(5.3), Inches(0.5), "豆包 / DeepSeek / 阿里千問 / 百度AI / 元寶 / Kimi", size=12.5, color=C["WHITE"])
    txt(s, Inches(7.05), Inches(5.8), Inches(5.3), Inches(0.5), "海外版 5（輔）：", size=13, color=C["WHITE"], bold=True)
    txt(s, Inches(7.05), Inches(6.2), Inches(5.3), Inches(0.5), "ChatGPT / Perplexity / Claude / Gemini / Copilot", size=12.5, color=C["WHITE"])

    # ================================================================
    # Slide 4 — 行業 AI 搜索現狀
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "行業 AI 搜索現狀", f"{PRODUCT_TYPE}零售的 GEO 機會（利基定位）")
    x = Inches(0.8)
    for big, label, col, src in STAT_CARDS:
        rect(s, x, Inches(1.6), Inches(3.7), Inches(3.5), fill=C["WHITE"], line=CM[col], line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x, Inches(2.0), Inches(3.7), Inches(1.3), big, size=46, color=CM[col], bold=True, align=PP_ALIGN.CENTER, font=FNUM)
        txt(s, x, Inches(3.4), Inches(3.7), Inches(0.9), label, size=14, color=C["INK"], align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.15), Inches(4.55), Inches(3.4), Inches(0.5), src, size=9.5, color=C["GRAY"], align=PP_ALIGN.CENTER)
        x += Inches(3.97)
    txt(s, Inches(0.8), Inches(5.5), Inches(11.6), Inches(1.3), INDUSTRY_NOTE, size=13.5, color=C["INK"])

    # ================================================================
    # Slide 5 — AI 平臺曝光度（客戶核心關切）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "AI 平臺曝光度 · 當前體檢與平台影響力", "客戶核心關切 · 國內 6 平臺視角（診斷模型估算）")
    rect(s, Inches(0.8), Inches(1.6), Inches(5.4), Inches(5.0), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.8), Inches(1.9), Inches(5.4), Inches(0.5), "國內 6 平臺平均被引用率", size=16, color=C["CYAN"], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(2.4), Inches(5.4), Inches(1.3), str(CITE_RATE), size=72, color=C["WHITE"], bold=True, align=PP_ALIGN.CENTER, font=FNUM)
    txt(s, Inches(0.8), Inches(3.75), Inches(5.4), Inches(0.5), f"{CITED} / {TOTAL_SCEN} 次模擬收錄場景中出現", size=13, color=C["WHITE"], align=PP_ALIGN.CENTER)
    rect(s, Inches(1.1), Inches(4.45), Inches(4.8), Inches(1.5), fill=C["BLUE"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(1.1), Inches(4.6), Inches(4.8), Inches(0.55), f"曝光缺口 {GAP_PCT}%", size=22, color=C["WHITE"], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.1), Inches(5.15), Inches(4.8), Inches(0.75), f"每 100 次相關提問，僅 {int(CITE_RATE*100)} 次見到本品牌；優秀線 0.70，尚差 {round(0.70-CITE_RATE,2)}", size=11.5, color=C["WHITE"], align=PP_ALIGN.CENTER)
    if VIS_DUAL:
        txt(s, Inches(0.85), Inches(6.0), Inches(5.4), Inches(0.5),
            f"雙指標（診斷模型估算）：品牌詞 {VIS_DUAL['brand_word']:.0%} · 品類詞 {VIS_DUAL['category_word']:.0%}（行業頭部≈{VIS_DUAL['industry_top']:.0%}）",
            size=10.5, color=C["CYAN"])
    txt(s, Inches(6.6), Inches(1.6), Inches(6.0), Inches(0.5), "平台影響力排名（對總被提及 %d 次的貢獻）" % CITED, size=15, color=C["BLUE"], bold=True)
    table(s, Inches(6.6), Inches(2.15), Inches(6.0), Inches(3.3), PLATFORM_INFLUENCE,
          [2.2, 1.8, 1.2, 1.8], fs=12.5, rh=0.5)
    txt(s, Inches(6.6), Inches(5.7), Inches(6.0), Inches(0.9), INFLUENCE_NOTE, size=11.5, color=C["INK"])

    # ================================================================
    # Slide 6 — AI 平臺常見問答圖譜
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "AI 平臺常見問答圖譜", "按用戶意圖聚類（微盟星启問答提取 · 國內版為主）")
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), QA_GRAPH,
          [1.4, 4.6, 2.2, 3.5], fs=12.5, rh=0.62)
    txt(s, Inches(0.8), Inches(6.75), Inches(11.6), Inches(0.35),
        f"來源：微盟星启問答提取（15 組高頻問句 × 國內 6 平臺檢索）。海外版問答以英文口徑規劃。", size=9.5, color=C["GRAY"])

    # ================================================================
    # Slide 7 — 高頻問題被提及次數（客戶核心關切）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "高頻問題被提及次數", f"15 組問句 × 國內 6 平臺 = {TOTAL_SCEN} 場景，共 {CITED} 次被提及（診斷模型估算）")
    dataB = [["意圖", "問句數", "被提及次數", "提及率", "主要觸發平臺"]] + FREQ_ROWS
    table(s, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.6), dataB,
          [2.4, 1.6, 2.2, 1.6, 3.9], fs=12.5, rh=0.62)
    txt(s, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.5), FREQ_HL, size=12, color=C["INK"])

    # ================================================================
    # Slide 8 — 常見引用文章 / 信源清單
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "常見引用文章 / 信源清單", "信源 typology 與覆蓋缺口（來源：基建診斷檢索）")
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.6), SOURCE_TABLE,
          [2.0, 5.1, 1.6, 3.0], fs=12, rh=0.7)
    _src_cite = ""
    if SOURCE_CITE:
        _src_cite = " · 信源引用次數榜（Top）：" + " / ".join(f"{n} {c}" for n, c in SOURCE_CITE[:5])
    txt(s, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.5),
        SOURCE_FOOTNOTE + _src_cite, size=11.5, color=C["INK"])

    # ================================================================
    # Slide 9 — 引用偏好框架（7 維度）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "從 AI 喜好入手 · 引用偏好框架（7 維度）", "可被 AI 采信度評分：60 / 100")
    data = [
        ["偏好維度", "對客戶的含義", "現狀", "缺口"],
        ["權威信源", "需官方 / 高權重站點背書", "✓ 具備", "—"],
        ["知識圖譜", "品牌—產品—門市結構化", "△ 部分", "需建立"],
        ["結構化 FAQ", "問答成對、命中答案抽取", "✗ 不足", "重點補"],
        ["數據與證據", "用事實替代形容詞", "△ 部分", "補評測"],
        ["新鮮度", "定期更新、有發布時間", "△ 部分", "建節奏"],
        ["多語言", "繁 / 簡 / 英分流客群", "✗ 不足", "補簡中"],
        ["多模態", "圖 / 影片 / 規格表", "△ 部分", "補多媒"],
    ]
    table(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(5.0), data,
          [2.4, 4.6, 1.8, 2.9], fs=12, rh=0.55)
    txt(s, Inches(0.8), Inches(6.7), Inches(11.6), Inches(0.35),
        "✓ 具備 / △ 部分 / ✗ 不足。結構化 FAQ 與簡中多語為最優先補強項，直接對應微盟星启「內容創作與改造優化」模塊。",
        size=10.5, color=C["GRAY"])

    # ================================================================
    # Slide 10 — AIVO 評分卡 + 雷達圖
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "AIVO 四維評分卡", f"總分 {AIVO_TOTAL} · {AIVO_RATE}（國內 6 平臺視角）")
    chart_data = CategoryChartData()
    chart_data.categories = ["AI可見度", "基建完善度", "競爭優勢", "輿情健康度"]
    chart_data.add_series("本品牌", (AIVO_VIS, AIVO_INFRA, AIVO_COMP, AIVO_SENT))
    chart_data.add_series("行業基準", tuple(AIVO_BENCH))
    gf = s.shapes.add_chart(XL_CHART_TYPE.RADAR, Inches(0.7), Inches(1.6), Inches(6.2), Inches(5.0), chart_data)
    chart = gf.chart
    chart.has_title = False; chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False; chart.legend.font.size = Pt(11)
    plot = chart.plots[0]; plot.has_data_labels = False
    chart.series[0].format.line.color.rgb = _cs0
    chart.series[0].format.line.width = Pt(2.25)
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = _cs0
    chart.series[1].format.line.color.rgb = _cs1
    chart.series[1].format.line.width = Pt(2)
    chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = _cs1
    aivo_notes = [
        ("AI 搜索可見度", str(AIVO_VIS), "國內6平臺 %.2f" % CITE_RATE, C["AMBER"]),
        ("基建完善度", str(AIVO_INFRA), "官網68+自媒體8+權威12", C["BLUE"]),
        ("競爭優勢", str(AIVO_COMP), "利基定位，落後頭部", C["BLUE"]),
        ("輿情健康度", str(AIVO_SENT), "負面率 ~2%（行業層級）", C["GREEN"]),
    ]
    y = Inches(1.7)
    for name, sc, note, col in aivo_notes:
        rect(s, Inches(7.2), y, Inches(5.3), Inches(1.12), fill=C["WHITE"], line=col, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, Inches(7.2), y, Inches(0.14), Inches(1.12), fill=col)
        txt(s, Inches(7.5), y + Inches(0.12), Inches(3.4), Inches(0.5), name, size=14, color=C["INK"], bold=True)
        txt(s, Inches(7.5), y + Inches(0.55), Inches(3.6), Inches(0.5), note, size=11, color=C["GRAY"])
        txt(s, Inches(11.4), y + Inches(0.18), Inches(1.0), Inches(0.8), sc, size=34, color=col, bold=True, align=PP_ALIGN.CENTER, font=FNUM)
        y += Inches(1.22)
    txt(s, Inches(7.2), Inches(6.7), Inches(5.3), Inches(0.4),
        "評級：≥90優 / ≥75良 / ≥60一般 / <60較差", size=10, color=C["GRAY"])

    # ================================================================
    # Slide 11 — 品牌基建診斷
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "品牌基建診斷（INFRA_EVAL）", "可被 AI 采信的內容資產現狀")
    x = Inches(0.8); cw = Inches(3.8)
    for title, sc, body, col in INFRA_CARDS:
        rect(s, x, Inches(1.6), cw, Inches(4.6), fill=C["WHITE"], line=CM[col], line_w=1.2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, x, Inches(1.6), cw, Inches(0.18), fill=CM[col], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, x + Inches(0.25), Inches(1.95), cw - Inches(0.5), Inches(0.5), title, size=18, color=CM[col], bold=True)
        txt(s, x + Inches(0.25), Inches(2.5), Inches(1.6), Inches(1.0), sc, size=44, color=CM[col], bold=True, font=FNUM)
        txt(s, x + Inches(1.9), Inches(2.6), cw - Inches(2.1), Inches(0.8), "項 / 條", size=13, color=C["GRAY"])
        txt(s, x + Inches(0.25), Inches(3.7), cw - Inches(0.5), Inches(2.3), body, size=13, color=C["INK"])
        x += cw + Inches(0.35)
    txt(s, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.4),
        "基建總評：官方站 + 權威媒體已具基礎，但「國內自媒體缺位」與「可被 AI 引用的簡中結構化深度」不足——正是微盟星启內容改造的切入點。",
        size=11.5, color=C["INK"])

    # ================================================================
    # Slide 12 — 輿情風險監控
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "輿情風險監控（SENTIMENT）", f"健康度 {AIVO_SENT} · 低風險 · 趨勢穩定")
    rect(s, Inches(0.8), Inches(1.6), Inches(4.0), Inches(2.4), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(0.8), Inches(1.8), Inches(4.0), Inches(0.5), "輿情健康度", size=16, color=C["CYAN"], bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(0.8), Inches(2.25), Inches(4.0), Inches(1.0), str(AIVO_SENT), size=60, color=C["WHITE"], bold=True, align=PP_ALIGN.CENTER, font=FNUM)
    txt(s, Inches(0.8), Inches(3.35), Inches(4.0), Inches(0.5), "負面率 ~2%（3/150）· 低風險", size=12.5, color=C["WHITE"], align=PP_ALIGN.CENTER)
    rect(s, Inches(5.1), Inches(1.6), Inches(7.4), Inches(5.0), fill=C["WHITE"], line=C["BLUE"], line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(5.35), Inches(1.8), Inches(7.0), Inches(0.5), "監測發現與風險清單", size=16, color=C["BLUE"], bold=True)
    bullets(s, Inches(5.35), Inches(2.4), Inches(7.0), Inches(2.0), SENTIMENT_BULLETS, size=12.5, gap=7)
    txt(s, Inches(5.35), Inches(4.6), Inches(7.0), Inches(0.5), "持續監控方案（微盟星启）", size=16, color=C["BLUE"], bold=True)
    bullets(s, Inches(5.35), Inches(5.2), Inches(7.0), Inches(1.3), SENTIMENT_ACTIONS, size=12.5, color=C["INK"], gap=6)

    # ================================================================
    # Slide 13 — 橫縱分析·企業現狀（縱向）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "橫縱分析 · 企業現狀（縱向）", "客觀事實 · 非主觀判斷")
    data = [["維度", "客觀事實（來源 / 證據）"]] + CURRENT_FACTS
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.0), data, [2.6, 9.1], fs=12.5, rh=0.72)

    # ================================================================
    # Slide 14 — 橫縱分析·行業與競品（橫向）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "橫縱分析 · 行業與競品（橫向）", "常見問答覆蓋率 / 被引用形式 / 信源 對比")
    eco_header = [["品牌", "AI 生態位置", "被引用形式", "主要信源"]] + ECO_MAP
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2), eco_header, [3.4, 3.0, 3.0, 2.3], fs=12, rh=0.42)
    _rank_txt = f"競品同題對比（國內 6 平臺 × {TOTAL_SCEN // 6} 問總提及）"
    if RANK_POOL:
        _rank_txt += f" · 競品池規模 {RANK_POOL['total']} 個品牌，本品牌排 #{RANK_POOL['rank']}"
    txt(s, Inches(0.8), Inches(4.25), Inches(11.6), Inches(0.35), _rank_txt, size=14, color=C["BLUE"], bold=True)
    data2 = [["品牌", "總提及", "評價"]]
    for name, cnt, note in COMPETITORS:
        data2.append([name, cnt, note])
    data2.append([BRAND_CN + "（本）", str(CITED), "長尾，利基定位"])
    table(s, Inches(0.8), Inches(4.65), Inches(11.7), Inches(2.2), data2, [4.4, 2.6, 4.7], fs=12, rh=0.42)

    # ================================================================
    # Slide 15 — 對手在 AI 平臺的做法（客戶核心關切）
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "對手在 AI 平臺的做法 · 競品基準", "他們如何佔住 AI 回答 · 本品牌可借鏡之處")
    rect(s, Inches(0.8), Inches(1.6), Inches(6.6), Inches(5.0), fill=C["WHITE"], line=C["BLUE"], line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(1.05), Inches(1.8), Inches(6.1), Inches(0.5), "對手在 AI 平臺的佔位套路", size=16, color=C["BLUE"], bold=True)
    bullets(s, Inches(1.05), Inches(2.4), Inches(6.1), Inches(4.0), [
        f"{c[0]}：總提及 {c[1]}，{c[2]}。" for c in COMPETITORS
    ], size=12.5, gap=9)
    rect(s, Inches(7.6), Inches(1.6), Inches(5.0), Inches(5.0), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(7.85), Inches(1.8), Inches(4.5), Inches(0.5), "本品牌可借鏡的三件套", size=16, color=C["CYAN"], bold=True)
    bullets(s, Inches(7.85), Inches(2.4), Inches(4.5), Inches(4.0), COMPETITOR_TACTICS,
            color=C["WHITE"], size=13, gap=10)
    txt(s, Inches(7.85), Inches(5.7), Inches(4.5), Inches(0.8), COMPETITOR_HL, size=11.5, color=C["WHITE"])

    # ================================================================
    # Slide 16 — GEO 優化問題清單
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "GEO 優化問題清單（數據驅動）", "對應常見問答 · 缺口類型 · 優先級")
    data = [["#", "對應常見問答", "缺口類型", "證據", "優先級"]] + GAP_ROWS
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), data, [0.6, 3.4, 2.3, 3.6, 1.0], fs=12, rh=0.62)

    # ================================================================
    # Slide 17 — 微盟星启 GEO 能力映射
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "微盟星启 GEO 能力映射", "問題 → 模塊 → 指標（數據追蹤閉環）")
    data = [["優化問題", "對應模塊", "追蹤指標", "優先級"]] + CAPABILITY_MAP
    table(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.6), data, [2.7, 3.5, 3.7, 1.0], fs=11.5, rh=0.6)
    txt(s, Inches(0.8), Inches(6.45), Inches(11.6), Inches(0.4),
        "四大能力模塊：AI 可見性監測 · 品牌輿情指數監測 · 內容創作與改造優化 · 智能媒體匹配與發布。",
        size=10.5, color=C["GRAY"])

    # ================================================================
    # Slide 18 — 數據追蹤方案
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "數據追蹤方案", "話題詞 / 平臺雙版本 / 媒體矩陣 / 基線目標")
    rect(s, Inches(0.8), Inches(1.55), Inches(5.6), Inches(2.0), fill=C["WHITE"], line=C["BLUE"], line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(1.0), Inches(1.7), Inches(5.2), Inches(0.4), "核心話題詞（來自問答圖譜）", size=14, color=C["BLUE"], bold=True)
    bullets(s, Inches(1.0), Inches(2.2), Inches(5.2), Inches(1.3), TOPIC_WORDS, size=12, gap=5)
    rect(s, Inches(6.6), Inches(1.55), Inches(5.9), Inches(2.0), fill=C["NAVY"], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(s, Inches(6.8), Inches(1.7), Inches(5.5), Inches(0.4), "目標平臺雙版本", size=14, color=C["CYAN"], bold=True)
    txt(s, Inches(6.8), Inches(2.2), Inches(5.5), Inches(0.5), "國內版 6（主）：豆包/DeepSeek/阿里千問/百度AI/元寶/Kimi", size=12, color=C["WHITE"])
    txt(s, Inches(6.8), Inches(2.75), Inches(5.5), Inches(0.5), "海外版 5（輔）：ChatGPT/Perplexity/Claude/Gemini/Copilot", size=12, color=C["WHITE"])
    txt(s, Inches(0.8), Inches(3.75), Inches(11.6), Inches(0.4), "基線 → 目標（可被引用率 / 覆蓋）", size=14, color=C["BLUE"], bold=True)
    table(s, Inches(0.8), Inches(4.05), Inches(11.7), Inches(2.4), KPI_ROWS, [4.4, 2.6, 2.6, 2.6], fs=12, rh=0.46)
    _comp_note = "海外版目標以 Phase1 診斷評估建立基線後滾動設定（本報告為規劃視角，不預設百分比）。"
    if COMPLIANCE.get("applicable"):
        _comp_note += " · 合規紅線：受限行業內容僅作成分/正品/品質/選購科普，不作功效聲稱（詳見話題詞方案）。"
    txt(s, Inches(0.8), Inches(6.6), Inches(11.6), Inches(0.35), _comp_note, size=9.5, color=C["GRAY"])

    # ================================================================
    # Slide 19 — 實施節奏與追蹤指標
    # ================================================================
    s = slide(); bg(s, C["CLOUD"])
    header(s, "實施節奏與追蹤指標", "對齊微盟星启五步標準服務流程")
    steps = [("1", "診斷評估"), ("2", "優化策略定制"), ("3", "執行上線"), ("4", "效果監控"), ("5", "持續優化")]
    n = len(steps); gap = Inches(0.3)
    cw = (SW - Inches(1.6) - gap * (n - 1)) / n
    x = Inches(0.8); y = Inches(1.7)
    for num, title in steps:
        rect(s, x, y, cw, Inches(2.3), fill=C["WHITE"], line=C["BLUE"], line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + cw/2 - Inches(0.42), y + Inches(0.22), Inches(0.84), Inches(0.84))
        circ.fill.solid(); circ.fill.fore_color.rgb = C["BLUE"]; circ.line.fill.background(); circ.shadow.inherit = False
        tf = circ.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = C["WHITE"]; r.font.name = FONT
        txt(s, x, y + Inches(1.2), cw, Inches(0.5), title, size=14, color=C["INK"], bold=True, align=PP_ALIGN.CENTER)
        if num != "5":
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + cw + Inches(0.01), y + Inches(0.55), gap - Inches(0.02), Inches(0.4))
            ar.fill.solid(); ar.fill.fore_color.rgb = C["CYAN"]; ar.line.fill.background(); ar.shadow.inherit = False
        x += cw + gap
    txt(s, Inches(0.8), Inches(4.4), Inches(11.6), Inches(0.4), "追蹤指標與節奏", size=14, color=C["BLUE"], bold=True)
    data = [["階段", "重點動作", "追蹤指標", "頻率"]] + ROADMAP_TABLE
    table(s, Inches(0.8), Inches(4.85), Inches(11.7), Inches(2.0), data, [3.0, 4.4, 3.0, 1.3], fs=11.5, rh=0.5)

    # ================================================================
    # Slide 20 — 結尾
    # ================================================================
    s = slide(); bg(s, C["NAVY"])
    rect(s, 0, SH - Inches(1.85), SW, Inches(1.85), fill=C["BLUE"])
    rect(s, 0, 0, Inches(0.25), SH, fill=C["CYAN"])
    txt(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(1.2), "讓品牌，在 AI 推薦中被看見", size=40, color=C["WHITE"], bold=True)
    txt(s, Inches(0.9), Inches(2.75), Inches(11.5), Inches(0.5), BRAND_CN + " · 微盟星启 GEO 優化方案", size=16, color=C["CYAN"])
    txt(s, Inches(0.9), Inches(3.65), Inches(11.5), Inches(0.45), CONTACT[0] if len(CONTACT) > 0 else "", size=15, color=C["WHITE"], bold=True)
    txt(s, Inches(0.9), Inches(4.2), Inches(11.5), Inches(0.45), CONTACT[1] if len(CONTACT) > 1 else "", size=13, color=C["WHITE"])
    txt(s, Inches(0.9), Inches(4.65), Inches(11.5), Inches(0.45), CONTACT[2] if len(CONTACT) > 2 else "", size=13, color=C["WHITE"])
    txt(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.4),
        "微盟星启GEO运营团队 · 由微盟星启提供", size=11, color=C["WHITE"], align=PP_ALIGN.CENTER)

    prs.save(output_path)
    return output_path
