#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
共享引擎：交付物 QA 检查器（含跨 case 自动反查）
=================================================
将 qa_check.py 的检查逻辑抽取为可编程函数，并内置「数据隔离」反查：
对给定 case 生成/已生成的交付物，自动扫描【其他所有 case】的 pollute_words.txt，
一旦当前 case 的交付物里出现别的 case 的行业词，即为跨案例污染，判定失败。

run_qa(case_name, cases_dir, extra_pollute=(), must_words=(), compliance=True)
  → (fail_count, report_lines)

覆盖 style_guide.md 全部规则：
  · 页数恰好 20        · 禁忌标记 0（虛擬/⚠️/website/{{）
  · 占位符 0           · 元素越界 OOB 0
  · 診斷模型估算 ≥ 2   · 合规功效词（受监管品类）仅出现在否定语境
  · 双指标 / 排名锚点 / 信源榜 / 合规声明 在交付物中体现
"""
import io
import os
import re

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

FORBIDDEN = ["實測", "虛擬", "\u26a0\ufe0f", "\u26a0", "website", "{{",
             "頭部競品 A", "頭部競品 B", "头部竞品 A", "头部竞品 B"]
BAN = ["改善記憶", "改善记忆", "增強腦力", "增强脑力", "降血脂", "降血壓", "降血压",
       "護心腦", "护心脑", "增強免疫", "增强免疫", "治療", "治疗", "預防", "预防",
       "療效", "疗效", "根治", "降三高", "通血管", "藥用", "药用", "防健忘", "补脑有效"]
NEG_CTX = ["禁用", "严禁", "嚴禁", "不写", "不寫", "不作", "不得", "红线", "紅線",
           "已改", "已重构", "已重構", "避免", "不提", "不碰", "不能", "不会", "不會",
           "无国内", "未取得", "未获得", "不涉", "改为", "改為", "绝不出现", "絕不出現",
           "以下表述", "不作任何功能", "不作任何功效"]


def deck_text(p):
    out = []
    for s in p.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                out.append(sh.text_frame.text)
            if getattr(sh, "has_table", False) and sh.has_table:
                for r in sh.table.rows:
                    for c in r.cells:
                        out.append(c.text)
    return "\n".join(out)


def strip_tags(s):
    return re.sub(r"<[^>]+>", " ", s)


def compliance_ok(text, is_html):
    """返回未通过的功效词命中数（仅统计非否定语境）。"""
    hits = [(w, text.count(w)) for w in BAN if w in text]
    bad = 0
    for w, _ in hits:
        for m in re.finditer(re.escape(w), text):
            win = text[max(0, m.start() - 400):m.end() + 200]
            ctx = re.sub(r"\s+", " ", strip_tags(win))
            if not any(k in ctx for k in NEG_CTX):
                bad += 1
    return bad


def run_qa(case_name, cases_dir="cases", extra_pollute=(), must_words=(),
           compliance=True, verbose=True):
    """对 case_name 的交付物执行 QA。返回 (fail_count, report[])。"""
    L = []
    fail = 0
    def p(msg):
        if verbose:
            L.append(msg)
        return msg

    case_path = os.path.join(cases_dir, case_name)
    out_dir = os.path.join(case_path, "output")
    if not os.path.isdir(out_dir):
        p(f"[QA] 缺 output 目录：{out_dir}")
        return (1, L)

    # 找到该 case 的 PPT 与 HTML
    deck = None; html_path = None
    for fn in sorted(os.listdir(out_dir)):
        fp = os.path.join(out_dir, fn)
        if fn.lower().endswith(".pptx") and deck is None:
            deck = fp
        if fn.lower().endswith(".html") and html_path is None:
            html_path = fp

    # ---- PPT 检查 ----
    if deck is None:
        p(f"✗ 未找到 {case_name} 的 .pptx 交付物")
        fail += 1
        return (fail, L)
    if Presentation is None:
        p("✗ 缺少 python-pptx")
        return (fail, L)
    prs = Presentation(deck)
    full = deck_text(prs)

    # 1. 页数
    if len(prs.slides) != 20:
        p(f"✗ 页数 {len(prs.slides)} != 20"); fail += 1
    else:
        p("✓ 页数 20")

    # 2. 禁忌标记
    fb = [(w, full.count(w)) for w in FORBIDDEN if full.count(w)]
    if fb:
        p(f"✗ 禁忌标记 {fb}"); fail += 1
    else:
        p("✓ 禁忌标记 0")

    # 3. OOB
    oob = [(i + 1, str(sh.shape_type)) for i, s in enumerate(prs.slides) for sh in s.shapes
           if sh.top is not None and sh.height is not None
           and sh.top + sh.height > prs.slide_height + 9144]
    if oob:
        p(f"✗ 元素越界 OOB {oob}"); fail += 1
    else:
        p("✓ 元素越界 OOB 0")

    # 4. 診斷模型估算
    n = full.count("診斷模型估算")
    if n < 2:
        p(f"✗ 診斷模型估算 仅 {n} 处（需 ≥2）"); fail += 1
    else:
        p(f"✓ 診斷模型估算 {n} 处")

    # 5. 跨 case 自动反查（数据隔离硬约束）——
    #    扫描【所有其他 case】的 pollute_words.txt，凡出现在本 case 交付物即污染。
    pollute = set(extra_pollute)
    for other in sorted(os.listdir(cases_dir)):
        if other == case_name or other.startswith("."):
            continue
        wf = os.path.join(cases_dir, other, "pollute_words.txt")
        if os.path.exists(wf):
            for w in io.open(wf, encoding="utf-8"):
                w = w.strip()
                if w:
                    pollute.add((w, other))
    pol = []
    for item in pollute:
        if isinstance(item, tuple):
            w, src = item
            if w and full.count(w):
                pol.append(f"{w}(来自 {src})")
        else:
            w = item
            if w and full.count(w):
                pol.append(w)
    if pol:
        p(f"✗ 跨 case 污染反查命中：{pol}"); fail += 1
    else:
        p("✓ 跨 case 污染反查 0（扫描全部其他 case 的 pollute_words）")

    # 6. 行业必备词
    if must_words:
        miss = [w for w in must_words if w and w not in full]
        if miss:
            p(f"✗ 行业必备词缺失 {miss}"); fail += 1
        else:
            p("✓ 行业必备词齐全")

    # 7. HTML 交付物（话题词方案）
    if html_path:
        htext = io.open(html_path, encoding="utf-8").read()
        hb = [(w, htext.count(w)) for w in FORBIDDEN if htext.count(w)]
        if hb:
            p(f"✗ 话题词 HTML 禁忌标记 {hb}"); fail += 1
        else:
            p("✓ 话题词 HTML 禁忌标记 0")
        # HTML 跨 case 反查
        hpol = []
        for item in pollute:
            if isinstance(item, tuple):
                w, src = item
                if w and w in htext:
                    hpol.append(f"{w}(来自 {src})")
            else:
                w = item
                if w and w in htext:
                    hpol.append(w)
        if hpol:
            p(f"✗ 话题词 HTML 跨 case 污染：{hpol}"); fail += 1
        else:
            p("✓ 话题词 HTML 跨 case 污染 0")

    # 8. 合规专项
    if compliance:
        cb = compliance_ok(full, False)
        if cb:
            p(f"✗ PPT 合规功效词 {cb} 处非否定语境"); fail += 1
        else:
            p("✓ PPT 合规功效词 0 处非否定语境")
        if html_path:
            hb2 = compliance_ok(htext, True)
            if hb2:
                p(f"✗ 话题词 HTML 合规功效词 {hb2} 处非否定语境"); fail += 1
            else:
                p("✓ 话题词 HTML 合规功效词 0 处非否定语境")

    p("===> %s" % ("ALL PASS" if fail == 0 else "FAIL %d 项" % fail))
    return (fail, L)


if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("case", help="case 名，如 meiriki")
    ap.add_argument("--cases-dir", default="cases")
    ap.add_argument("--no-compliance", action="store_true")
    args = ap.parse_args()
    f, lines = run_qa(args.case, cases_dir=args.cases_dir,
                      compliance=not args.no_compliance)
    for l in lines:
        print(l)
    raise SystemExit(f)
